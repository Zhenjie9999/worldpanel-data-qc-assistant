from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


NUMBER_PATTERN = re.compile(r"[-+]?\d+(?:,\d{3})*(?:\.\d+)?\s*%?")


def _extract_numbers(text: str, context: str, location: str, file_name: str) -> list[dict]:
    results = []
    for match in NUMBER_PATTERN.finditer(text or ""):
        raw = match.group(0).strip()
        results.append(
            {
                "value": float(raw.replace(",", "").replace("%", "")),
                "raw": raw,
                "is_percent": "%" in raw,
                "context": context,
                "location": location,
                "file_name": file_name,
            }
        )
    return results


def parse_pptx(path: Path) -> dict:
    prs = Presentation(path)
    pages, all_numbers, texts = [], [], []
    for page_number, slide in enumerate(prs.slides, 1):
        page_texts, page_numbers = [], []
        structured_shapes = image_shapes = grouped_shapes = chart_shapes = 0
        for shape_index, shape in enumerate(slide.shapes, 1):
            location = f"Slide {page_number} / Shape {shape_index}"
            text = getattr(shape, "text", "") or ""
            if text.strip():
                cleaned = " ".join(text.split())
                page_texts.append(cleaned)
                texts.append({"text": cleaned, "location": location})
                page_numbers.extend(_extract_numbers(cleaned, cleaned, location, path.name))
            if getattr(shape, "has_table", False):
                structured_shapes += 1
                for row in shape.table.rows:
                    for cell in row.cells:
                        cleaned = " ".join(cell.text.split())
                        page_numbers.extend(_extract_numbers(cleaned, cleaned, location, path.name))
            if getattr(shape, "has_chart", False):
                structured_shapes += 1
                chart_shapes += 1
                for series in shape.chart.series:
                    context = f"{series.name} {' '.join(page_texts[:2])}".strip()
                    for value in series.values:
                        if isinstance(value, (int, float)):
                            page_numbers.append(
                                {
                                    "value": float(value),
                                    "raw": str(value),
                                    "is_percent": False,
                                    "context": context,
                                    "location": location,
                                    "file_name": path.name,
                                }
                            )
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_shapes += 1
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                grouped_shapes += 1
        review_required = image_shapes > 0 or grouped_shapes > 0 or chart_shapes > 0
        coverage = 70 if review_required else 100
        if structured_shapes and review_required:
            coverage = 85
        if chart_shapes and not image_shapes and not grouped_shapes:
            coverage = 95
        pages.append(
            {
                "page": page_number,
                "text": " | ".join(page_texts),
                "numbers": page_numbers,
                "coverage_percent": coverage,
                "numbers_found": len(page_numbers),
                "low_confidence_count": image_shapes + grouped_shapes,
                "review_required": review_required,
                "detail": (
                    f"{image_shapes} image region(s); {grouped_shapes} grouped region(s); "
                    f"{chart_shapes} chart(s); {structured_shapes} structured object(s)"
                ),
            }
        )
        all_numbers.extend(page_numbers)
    return {
        "file_name": path.name,
        "file_type": "pptx",
        "numbers": all_numbers,
        "texts": texts,
        "pages": pages,
        "warnings": [],
    }
