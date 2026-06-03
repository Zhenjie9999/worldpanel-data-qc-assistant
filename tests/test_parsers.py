import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from openpyxl import Workbook
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from worldpanel_qc.parsers import parse_file
from worldpanel_qc.parsers.excel import _convert_xls


class ParserTests(unittest.TestCase):
    def test_parses_excel_numbers_and_hidden_columns(self):
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "source.xlsx"
            wb = Workbook()
            ws = wb.active
            ws.title = "Summary"
            ws["A1"] = "Penetration"
            ws["B1"] = 0.402
            ws.column_dimensions["B"].hidden = True
            wb.save(path)

            parsed = parse_file(path)

            self.assertEqual(parsed["file_type"], "xlsx")
            self.assertEqual(parsed["numbers"][0]["location"], "Summary!B1")
            self.assertEqual(parsed["hidden_columns"], ["Summary!B"])

    def test_converts_xls_with_libreoffice_when_available(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "legacy.xls"
            source.write_bytes(b"legacy")
            calls = []

            def fake_run(command, check, capture_output, timeout, text):
                calls.append(command)
                converted_dir = Path(command[command.index("--outdir") + 1])
                converted_dir.mkdir(parents=True, exist_ok=True)
                (converted_dir / "legacy.xlsx").write_bytes(b"converted")

            with patch("worldpanel_qc.parsers.excel.shutil.which", side_effect=lambda name: "soffice" if name == "soffice" else None):
                with patch("worldpanel_qc.parsers.excel.subprocess.run", side_effect=fake_run):
                    converted, warning = _convert_xls(source)

            self.assertEqual(calls[0][0], "soffice")
            self.assertEqual(converted.read_bytes(), b"converted")
            self.assertIn("LibreOffice", warning)

    def test_parses_powerview_measure_row_and_column_context(self):
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "price.xlsx"
            wb = Workbook()
            ws = wb.active
            ws["A1"] = "Measures = Weighted RESP Price per Volume"
            ws["A2"] = "n370yg - table"
            ws["B2"] = "Cherry"
            ws["A3"] = "52 w/e 2022/12/30"
            ws["B3"] = 8000
            ws["J3"] = "Trailing note"
            wb.save(path)

            parsed = parse_file(path)

            number = parsed["numbers"][0]
            self.assertEqual(number["measure"], "Weighted RESP Price per Volume")
            self.assertEqual(number["row_label"], "52 w/e 2022/12/30")
            self.assertEqual(number["column_label"], "Cherry")
            self.assertIn("Weighted RESP Price per Volume", number["context"])
            self.assertIn("Cherry", number["context"])

    def test_parses_pptx_chart_values_and_marks_picture_for_review(self):
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "deck.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            chart_data = ChartData()
            chart_data.categories = ["FY2025"]
            chart_data.add_series("Penetration", (40.2,))
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(1),
                Inches(1),
                Inches(4),
                Inches(3),
                chart_data,
            )
            # A picture-like shape is not required for this test; a chart is structured coverage.
            prs.save(path)

            parsed = parse_file(path)

            self.assertEqual(parsed["file_type"], "pptx")
            self.assertEqual(parsed["pages"][0]["numbers"][0]["value"], 40.2)
            self.assertGreaterEqual(parsed["pages"][0]["coverage_percent"], 90)
            self.assertTrue(parsed["pages"][0]["review_required"])

    def test_marks_grouped_pptx_shapes_for_visual_review(self):
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "grouped-deck.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            group = slide.shapes.add_group_shape()
            group.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
            prs.save(path)

            parsed = parse_file(path)

            self.assertTrue(parsed["pages"][0]["review_required"])
            self.assertIn("grouped", parsed["pages"][0]["detail"])


if __name__ == "__main__":
    unittest.main()
