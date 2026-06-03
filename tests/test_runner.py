import tempfile
import unittest
from pathlib import Path

from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from worldpanel_qc.qc.runner import run_qc


class RunnerTests(unittest.TestCase):
    def test_runner_reports_reading_local_rules_and_report_stages(self):
        with tempfile.TemporaryDirectory() as tmp:
            excel = Path(tmp) / "source.xlsx"
            wb = Workbook()
            wb.active["A1"] = "Penetration"
            wb.active["B1"] = 0.402
            wb.save(excel)
            progress = []

            run_qc(
                [excel],
                project_rules=[],
                external_ai_enabled=False,
                progress_callback=lambda stage, percent, detail="": progress.append((stage, percent, detail)),
            )

            self.assertIn(("Reading files", 8, "Reading file 1 / 1"), progress)
            self.assertTrue(any(stage == "Running local rules" for stage, _, _ in progress))
            self.assertTrue(any(stage == "Preparing report" for stage, _, _ in progress))
    def test_runner_returns_placeholder_issue_and_coverage(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            excel = tmp_path / "source.xlsx"
            wb = Workbook()
            ws = wb.active
            ws["A1"] = "Penetration"
            ws["B1"] = 0.402
            wb.save(excel)

            pptx = tmp_path / "deck.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            box.text = "Report/Presentation Name"
            prs.save(pptx)

            result = run_qc([excel, pptx], project_rules=[], external_ai_enabled=False)

            self.assertTrue(any(issue["rule_id"] == "placeholder_text" for issue in result["issues"]))
            self.assertEqual(len(result["coverage"]), 1)

    def test_mapping_constraint_limits_cross_file_source_candidates(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            first = tmp_path / "first.xlsx"
            wb = Workbook()
            ws = wb.active
            ws.title = "Summary"
            ws["A1"] = "Penetration"
            ws["B1"] = 0.402
            wb.save(first)
            second = tmp_path / "second.xlsx"
            wb.save(second)

            pptx = tmp_path / "deck.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            box.text = "Penetration 40.2%"
            prs.save(pptx)

            result = run_qc(
                [first, second, pptx],
                project_rules=[],
                external_ai_enabled=False,
                mapping_constraints=[
                    {
                        "page_file_name": "deck.pptx",
                        "page": 1,
                        "source_file_name": "second.xlsx",
                        "sheet_name": "Summary",
                    }
                ],
            )

            self.assertTrue(result["matches"])
            self.assertTrue(all(candidate["file_name"] == "second.xlsx" for candidate in result["matches"][0]["candidates"]))

    def test_runner_keeps_visible_numbers_without_excel_candidate(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            excel = tmp_path / "source.xlsx"
            wb = Workbook()
            ws = wb.active
            ws["A1"] = "Penetration"
            ws["B1"] = 0.402
            wb.save(excel)

            pptx = tmp_path / "deck.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            box.text = "Share 30.0%"
            prs.save(pptx)

            result = run_qc([excel, pptx], project_rules=[], external_ai_enabled=False)

            self.assertEqual(len(result["matches"]), 1)
            self.assertEqual(result["matches"][0]["status"], "unmatched")
            self.assertEqual(result["matches"][0]["candidates"], [])

    def test_runner_adds_local_share_total_issue(self):
        with tempfile.TemporaryDirectory() as tmp:
            excel = Path(tmp) / "source.xlsx"
            wb = Workbook()
            ws = wb.active
            ws["A1"] = "Brand Share"
            ws["B1"] = 0.4
            ws["B1"].number_format = "0.0%"
            ws["C1"] = 0.3
            ws["C1"].number_format = "0.0%"
            ws["D1"] = 0.2
            ws["D1"].number_format = "0.0%"
            wb.save(excel)

            result = run_qc([excel], project_rules=[], external_ai_enabled=False)

            self.assertTrue(any(issue["rule_id"] == "local_share_total" for issue in result["issues"]))

    def test_runner_adds_powerview_price_outlier_issue(self):
        with tempfile.TemporaryDirectory() as tmp:
            excel = Path(tmp) / "price.xlsx"
            wb = Workbook()
            ws = wb.active
            ws["A1"] = "Measures = Weighted RESP Price per Volume"
            ws.append(["n370yg - table", "Cherry"])
            ws.append(["52 w/e 2021/12/31", 47.6439])
            ws.append(["52 w/e 2022/12/30", 8000])
            ws.append(["52 w/e 2023/12/29", 63.7041])
            ws.append(["52 w/e 2024/12/27", 68.6306])
            wb.save(excel)

            result = run_qc([excel], project_rules=[], external_ai_enabled=False)

            issue = next(issue for issue in result["issues"] if issue["rule_id"] == "local_outlier")
            self.assertEqual(issue["location"], "Sheet!B4")
            self.assertIn("8000.0", issue["description"])


if __name__ == "__main__":
    unittest.main()
